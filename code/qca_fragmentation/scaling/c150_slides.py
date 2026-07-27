"""
Build an editable PowerPoint deck for the C150 discussion.

    python -m qca_fragmentation.scaling.c150_slides

Writes reports/slides/C150_discussion.pptx.  Everything is native PowerPoint --
placeholder text boxes, real bullet lists, real tables, and the figures as
pictures -- so the deck can be restyled and rewritten without touching this
script.  Numbers are pulled from the analytics JSONs where possible so a rerun
after new data stays consistent; the speaker notes carry the provenance
(report section, file) for each slide.
"""

from __future__ import annotations

import argparse
import json
import os
from math import comb, log
from typing import List, Optional, Sequence, Tuple

from .. import c150
from ..results_io import REPO_ROOT, load_results, sizes_from_record

ANALYTICS = os.path.join(REPO_ROOT, "analytics")
FIGDIR = os.path.join(REPO_ROOT, "figures")
SLIDEDIR = os.path.join(REPO_ROOT, "reports", "slides")
OUT = os.path.join(SLIDEDIR, "C150_discussion.pptx")

TITLE_PT = 26          # the stock 44 pt wraps on a 16:9 slide; edit freely


def _load(name: str) -> Optional[dict]:
    p = os.path.join(ANALYTICS, name)
    if not os.path.exists(p):
        return None
    with open(p) as f:
        return json.load(f)


# --- slide helpers -----------------------------------------------------------

def _new_deck():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)      # 16:9
    prs.slide_height = Inches(7.5)
    return prs


def _widen(prs, s, body_top: float = 1.55, body_height: float = 5.5):
    """
    The stock template's placeholders are laid out for 4:3, so on a 16:9 slide
    titles wrap and bodies sit in a narrow column.  Stretch them to the slide.
    """
    from pptx.util import Inches

    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        ph.left = Inches(0.6)
        ph.width = Inches(12.1)
        if idx == 0:                       # title
            ph.top = Inches(0.3)
            ph.height = Inches(1.0)
            from pptx.util import Pt
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(TITLE_PT)
        elif ph.has_text_frame:            # body / subtitle
            ph.top = Inches(body_top)
            ph.height = Inches(body_height)
    return s


def _title_slide(prs, title: str, subtitle: str, notes: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    _widen(prs, s, body_top=3.4, body_height=2.2)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _bullets(prs, title: str, items: Sequence, notes: str = "",
             font_size: int = 18):
    """
    items: a bare string is a SUB-bullet (level 1, the common case here);
    a tuple (text, level) or (text, level, bold) sets the level explicitly, so
    headings are written as (text, 0, True).
    """
    from pptx.util import Pt

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    _widen(prs, s)
    tf = s.placeholders[1].text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, str):
            text, level, bold = it, 1, False
        elif len(it) == 2:
            (text, level), bold = it, False
        else:
            text, level, bold = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = str(text)
        p.level = int(level)
        for r in p.runs:
            r.font.size = Pt(font_size - 2 * int(level))
            r.font.bold = bool(bold)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _picture(prs, title: str, image: str, caption: str = "", notes: str = ""):
    from pptx.util import Inches, Pt

    s = prs.slides.add_slide(prs.slide_layouts[5])       # title only
    s.shapes.title.text = title
    _widen(prs, s)
    if os.path.exists(image):
        pic = s.shapes.add_picture(image, Inches(0.6), Inches(1.6),
                                   width=Inches(12.1))
        if pic.height > Inches(4.9):
            scale = Inches(4.9) / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
            pic.left = int((prs.slide_width - pic.width) / 2)
    if caption:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1),
                                  Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = caption
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(13)
            r.font.italic = True
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _table(prs, title: str, header: Sequence[str], rows: Sequence[Sequence],
           notes: str = "", col_widths: Optional[Sequence[float]] = None,
           font_size: int = 13, top: float = 1.7):
    from pptx.util import Inches, Pt

    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    _widen(prs, s)
    nrow, ncol = len(rows) + 1, len(header)
    height = min(5.2, 0.34 * nrow + 0.15)
    shape = s.shapes.add_table(nrow, ncol, Inches(0.6), Inches(top),
                               Inches(12.1), Inches(height))
    tbl = shape.table
    if col_widths:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(12.1 * w / total)
    for j, h in enumerate(header):
        c = tbl.cell(0, j)
        c.text = str(h)
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(font_size)
                r.font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# --- the deck ----------------------------------------------------------------

def build(out: str = OUT) -> str:
    prs = _new_deck()
    lv = _load("c150_levels.json")
    hsf = _load("c150_hsf_compare.json")
    ent = _load("c150_entanglement.json")

    # 1 ------------------------------------------------------------------
    _title_slide(
        prs,
        "C150: an exactly solved quantum cellular automaton",
        "…and exactly what the solution does not determine\n"
        "Krylov sector analysis, Tier 1 — report R8",
        "Framing for the discussion: C150 is the one rule in the 16-rule "
        "unitary family whose sector structure we can write down in closed "
        "form. That makes it the right place to ask what a sector does and "
        "does not tell you — which is the question that generalises.",
    )

    # 2 ------------------------------------------------------------------
    _bullets(prs, "The rule, and one naming note", [
        ("C150 = Wolfram 150 = local symbol table (I, V, V, I)", 0, True),
        "Hadamard fires at site i  ⟺  x(i−1) ⊕ x(i+1) = 1",
        "  — the firing condition is exactly classical ECA-90's update",
        ("Boundary: obc0, the missing neighbour reads 0", 0, True),
        ("Cycle: brick-wall — even sublattice first, then odd", 0, True),
        ("Same object as HSF rule 6 with gate H (checked from HSF/src/model.jl)", 0, True),
        "  — so the whole Julia eigendata set is a usable independent oracle",
        ("Why this rule", 0, True),
        "R2 saw binomial sector sizes; R1 uses it as the HSF regression fixture;",
        "and its fitted growth exponent 2.022 sat above the informational ceiling 2",
    ], notes="R8 §1. The 2.022 anomaly is resolved on slide 6 — it was a "
             "prefactor being read as a base.")

    # 3 ------------------------------------------------------------------
    _bullets(prs, "Five results", [
        ("1. Reduction — sector partition = a single-flip graph", 0, True),
        "for ANY unitary rule; the brick-wall ordering drops out",
        ("2. Solution — walls, hard-core hopping, closed form", 0, True),
        "|S_w| = C(N+1, w) for even w; the charge is the Ising domain-wall number",
        ("3. C150 is symmetry-resolved, NOT fragmented", 0, True),
        "sector count is linear in N — it is the family's null model",
        ("4. Interacting, but not chaotic", 0, True),
        "#distinct eigenvalues = 3^⌈N/2⌉ exactly — no RMT ensemble fits",
        ("5. Area law vs volume law is a sector statement", 0, True),
        "…but the sector fixes only the ceiling, not the value",
    ], notes="The through-line: kinematics is solved exactly; everything "
             "dynamical sits strictly below what kinematics allows.")

    # 4 ------------------------------------------------------------------
    _bullets(prs, "Result 1 — the sector partition is a flip graph", [
        ("Theorem. For a unitary rule the sectors are the connected components of",
         0, True),
        "x  ~  x ⊕ 2^i   for every site i where the Hadamard fires",
        ("Proof sketch", 0, True),
        "each half-layer is depth-1 ⇒ its image from |x⟩ is a subcube (R1 no-interference)",
        "x lies in its own image (⟨b|H|b⟩ ≠ 0) ⇒ composite image ⊇ both half-layer images",
        "inside a half-layer the V/I pattern is frozen ⇒ the subcube is single-flip connected",
        ("Two consequences", 0, True),
        "the brick-wall ORDER is irrelevant — the partition is a property of the symbol table",
        "cost drops from O(2^N·|succ|) to O(N·2^N) — this is what reaches N = 32",
        ("Validated: all 16 unitary rules × both boundaries × N = 3…12 — 320 units, 0 mismatches",
         0, True),
    ], notes="R8 §2. The validation is against the streamed engine, which "
             "knows nothing about the theorem.")

    # 5 ------------------------------------------------------------------
    _bullets(prs, "Result 2 — walls: the whole rule in one line", [
        ("Bond variables  b_j = x_j ⊕ x_(j+1),  j = 0…N,  x(−1) = x(N) = 0", 0, True),
        "x ↦ b is a BIJECTION onto the even-weight strings of length N+1",
        ("The elementary move", 0, True),
        "flipping x_i toggles bonds b_i, b_(i+1), and is allowed iff b_i ≠ b_(i+1)",
        "⇒ a HARD-CORE HOP of one domain wall.  Walls are never created or destroyed.",
        ("Therefore", 0, True),
        "conserved charge Q = Σ (1 − Z_j Z_(j+1))/2 — the open Ising domain-wall number",
        "each wall shell is exactly one sector (hard-core hopping is connected on a shell)",
        ("Zero charge violations over every state, N = 3…12, both boundary conditions",
         0, True),
    ], notes="R8 §3. This is the step that turns a numerical observation into "
             "a closed form. Everything downstream is bookkeeping on this picture.")

    # 6 ------------------------------------------------------------------
    _bullets(prs, "Result 3 — the closed form, and the growth law", [
        ("|S_w| = C(N+1, w) for even w     #sectors = ⌊(N+1)/2⌋ + 1     "
         "D_max = max over even w", 0, True),
        "matches every stored Tier-1a unit, N = 6…20, both boundary conditions",
        ("C150 is NOT fragmented", 0, True),
        "sector count is LINEAR in N and every sector is a full Q-eigenspace",
        "no Krylov structure beyond the U(1) charge — contrast W204's 2^N frozen singletons",
        "→ it is the wrong example to quote for Hilbert-space fragmentation",
        ("The 2.022 anomaly, resolved", 0, True),
        "D_max/2^N = 2√(2/π(N+1)) — base is EXACTLY 2, the decay is an N^(−1/2) prefactor",
        "plus a sawtooth at N ≡ 1 (mod 4), where the central binomial sits at odd w",
        "a pure-exponential fit reads prefactor + sawtooth as a base above the ceiling",
    ], notes="R8 §4. This also corrected two loose statements in R2: sector "
             "sizes are the EVEN-WEIGHT entries of Pascal's row N+1 (the "
             "'lower half' phrasing only coincides for even N), and D_max is "
             "the central binomial except at N ≡ 1 mod 4.")

    # 7 ------------------------------------------------------------------
    eng = load_results(150, "obc0")
    fr = c150.load_frontier()
    rows = []
    for N in (20, 21, 22, 28, 30, 31, 32):
        cf = c150.sector_sizes_obc0(N)
        e = eng.get(N)
        f = fr.get(f"{N}_obc0")
        et = "—"
        if e and sizes_from_record(e) == cf:
            et = f"{e['runtime']/3600:.1f} h" if e["runtime"] > 3600 else \
                 f"{e['runtime']:.0f} s"
        ft = "—"
        if f and f["closed_form_exact"]:
            ft = f"{f['runtime']:.1f} s" if f["runtime"] < 60 else \
                 f"{f['runtime']:.0f} s"
        rows.append([N, f"{1 << N:,}", len(cf), f"{cf[0]:,}", et, ft])
    _table(prs, "Numerical frontier — two independent routes",
           ["N", "2^N", "#sectors", "D_max", "streamed engine", "flip graph"],
           rows,
           notes="R8 §5. Left route uses exact Z[1/√2] amplitudes and no "
                 "result of this report; right route uses the flip reduction. "
                 "They agree at N=21,22 where they overlap, and the reduction "
                 "is 4×10^4 / 1×10^5 times faster there. N=32 = 4.29e9 basis "
                 "states, 21.0 GB peak — the uint32 and RAM ceiling; N=33 "
                 "would need a 64-bit label array (68.7 GB).",
           col_widths=[1, 2, 1.4, 2, 2.2, 2.2])

    # 8 ------------------------------------------------------------------
    _picture(prs, "Exact sector structure, and where the numerics stop",
             os.path.join(FIGDIR, "fig_c150_frontier.png"),
             "Left: the closed form with both computed frontiers marked. "
             "Right: D_max/2^N against 2√(2/π(N+1)) — base 2, N^(−1/2) prefactor, "
             "sawtooth at N ≡ 1 (mod 4).",
             notes="Talking point: the closed form is exact at every N, so the "
                   "frontier is a validation frontier, not a discovery frontier.")

    # 9 ------------------------------------------------------------------
    _bullets(prs, "What the sector does NOT determine", [
        ("A sector is the maximal reachable set from a basis state — nothing more",
         0, True),
        ("An exponentially large stationary space", 0, True),
        "dim Fix(U_w) = C(⌈N/2⌉, w/2),  dim Fix(U) = 2^⌈N/2⌉ = 2^(#even sites)",
        "Fix(U) = Fix(L_A) ∩ Fix(L_B); almost all of it is DARK SUPERPOSITIONS",
        ("Hence an inner sector is never a Krylov space", 0, True),
        "spectrum degenerate ⟺ |S_w| > N+1;  dim K(x)/|S_w| = 0.40 … 0.96",
        ("Monitored vs unmonitored disagree inside a sector", 0, True),
        "monitored: provably uniform on S_w (doubly stochastic, irreducible, aperiodic)",
        "unmonitored: diagonal ensemble occupies only 6.6 % … 50 % of the sector",
        ("Reflection: commutes for odd N; conjugates U to U⁻¹ for even N", 0, True),
    ], notes="R8 §6. This is the precise, rule-specific form of the general "
             "caveat in R1 §'What a sector is': the graph gives dim K(x) ≤ "
             "|S_w| and away from the extremal shells the inequality is strict.")

    # 10 -----------------------------------------------------------------
    _picture(prs, "The graph gives the container, not the filling",
             os.path.join(FIGDIR, "fig_c150_quantum.png"),
             "Left: Krylov and diagonal-ensemble fractions of a sector; the "
             "dashed line is the monitored answer, which is exactly uniform. "
             "Right: dim Fix(U) = 2^⌈N/2⌉.",
             notes="Open markers on the left are the extremal shells "
                   "|S_w| ≤ N+1 — the only ones with a nondegenerate spectrum, "
                   "hence the only ones a basis state's Krylov space fills.")

    # 11 -----------------------------------------------------------------
    _bullets(prs, "Result 4a — interacting, not free", [
        ("In the bond basis the gate at site i touches bonds i, i+1 only", 0, True),
        "identity when both bonds agree; on the one-wall subspace a Hadamard coin",
        "M(c) = (1/√2)[[−(−1)^c, 1], [1, (−1)^c]],  det M = −1,  c = wall parity to the left",
        "this bond circuit reproduces the engine to 1.1×10⁻¹⁶ — it IS C150",
        ("Two diagonal signs spoil gaussianity", 0, True),
        "det(H) = −1 but the gate acts by +1 on a doubly occupied pair → maximal "
        "nearest-neighbour wall interaction exp(iπ n n)",
        "the (−1)^c string sits on the DENSITY term, not the hop → not absorbed by Jordan–Wigner",
        ("Sharp test: a Gaussian circuit has additive many-body eigenphases", 0, True),
        "C150: 0.55 … 1.87 rad   remove one defect: still 0.43 … 1.10 rad",
        "remove BOTH: additive to 4×10⁻¹⁵ — so both signs are needed, and C150 is interacting",
    ], notes="R8 §7. Careful: interacting does NOT imply chaotic — see the "
             "next slide. That inference is exactly the one we had to retract.")

    # 12 -----------------------------------------------------------------
    dist = (lv or {}).get("distinct", [])
    rows = [[r["N"], f"{r['dim']:,}", f"{r['n_distinct']:,}",
             f"{r['closed_form']:,}", f"{r['distinct_fraction']:.3f}"]
            for r in dist if r["N"] >= 8]
    _table(prs, "Result 4b — the spectrum is exponentially degenerate",
           ["N", "2^N", "#distinct", "closed form 3^m / (3^m±1)/2", "fraction"],
           rows,
           notes="R8 §7.3. m = ⌈N/2⌉ = number of even sites. The law was "
                 "fitted on N ≤ 13 and confirmed OUT OF SAMPLE at N = 14 "
                 "(predicted 2187, got 2187). No tolerance to argue about: at "
                 "N=12 the gaps split into 3367 below 1e-12 and 728 above "
                 "1e-4 with five empty decades between, so the count is "
                 "identical for thresholds 1e-13 … 1e-5.",
           col_widths=[1, 2, 2, 3.2, 1.6], font_size=12)

    # 13 -----------------------------------------------------------------
    refs = (lv or {}).get("references", {})
    def _r(k):
        v = refs.get(k)
        return f"{v['r_tilde_mean']:.3f}" if v else "—"
    _bullets(prs, "Result 4c — no random-matrix ensemble fits", [
        ("β = 2 (GUE / CUE) excluded by symmetry", 0, True),
        "U is real orthogonal ⇒ complex conjugation is a time reversal, K²=1",
        "⇒ eigenvalues in mirror pairs; all statistics on θ ∈ (0, π) only",
        ("β = 1 (COE / Haar-O) excluded by the degeneracy law", 0, True),
        "no chaotic Floquet operator has 3^m distinct eigenvalues out of 2^N",
        ("Poisson excluded too", 0, True),
        "ONE universal set Θ_N per N: every shell's spectrum ⊆ Θ_N, and the "
        "largest shell realises ALL of it (small shells as little as 1.8 %)",
        f"largest shell N=16, w=8: ⟨r̃⟩ = 0.347, BELOW Poisson ({_r('poisson')}) — "
        "and superposition can only approach Poisson from above",
        ("Structural source: U = L_B·L_A is a product of two involutions", 0, True),
        "gives the ±θ pairing and explains Fix(U) = Fix(L_A) ∩ Fix(L_B)",
        ("⇒ C150 is INTERACTING BUT NOT CHAOTIC — the two are compatible", 0, True),
    ], notes=f"R8 §7.4–7.5. References generated through the identical "
             f"pipeline: Poisson {_r('poisson')}, 2×Haar-O {_r('orth_x2')}, "
             f"COE {_r('coe')}, Haar-O {_r('orth')}, CUE {_r('cue')}. "
             "Erratum worth mentioning: an earlier draft predicted the "
             "statistics 'will not look free', which read as 'will look "
             "chaotic' was wrong.")

    # 14 -----------------------------------------------------------------
    _picture(prs, "Level statistics against generated references",
             os.path.join(FIGDIR, "fig_c150_levels.png"),
             "Left: ⟨r̃⟩ per symmetry block with the reference ensembles as "
             "bands. Right: the unfolded spacing distribution.",
             notes="Three things had to be handled or the numbers are "
                   "meaningless: odd N needs reflection-parity resolution "
                   "(0.43 → 0.60 at N=11); the odd-N shells w and N+1−w have "
                   "identical spectra and are counted once; and fixed w with "
                   "growing N is a FEW-BODY problem that must go Poisson.")

    # 15 -----------------------------------------------------------------
    rows = []
    for r in (hsf or []):
        rows.append([r["N"], f"{r['n_states']:,}", r["n_sectors"],
                     f"{r['n_distinct']:,}", f"{r['distinct_closed_form']:,}",
                     "✓" if r["sizes_match_closed_form"] else "✗",
                     "✓" if r["dim_fix_matches"] else "✗",
                     f"{r['frac_states_in_degenerate_eigenspace']:.2f}"])
    _table(prs, "Independent check: the HSF Julia eigendata (rule 6 = 150)",
           ["N", "states", "sectors", "#distinct", "3^m law",
            "sizes", "dim Fix", "deg. states"], rows,
           notes="R8 §8. Same operator, verified by reading HSF/src/model.jl; "
                 "only a bit reversal of the basis and θ = −arg λ differ, "
                 "neither of which can touch a spectrum. N=14 = 2187 is the "
                 "out-of-sample prediction reproduced by a different language, "
                 "sparse assembly, basis ordering and sign convention. "
                 "CAVEAT for the stored entropies — see next slide.",
           col_widths=[1, 1.8, 1.4, 1.8, 1.6, 1.2, 1.4, 1.8], font_size=12)

    # 16 -----------------------------------------------------------------
    _bullets(prs, "A caveat for the stored entanglement entropies", [
        ("On a degenerate eigenvalue the eigensolver returns an ARBITRARY "
         "orthonormal basis", 0, True),
        "and the half-chain entropy is not constant on that basis",
        ("Measured", 0, True),
        "rotating inside one eigenspace moves the stored entropies by up to 0.27 nats",
        "the same test on a simple eigenvalue gives a spread of exactly 0",
        ("For this rule that is the typical case, not an edge case", 0, True),
        "fraction of stored eigenstates in a degenerate eigenspace: 0.43 (N=8) → 0.74 (N=14)",
        "multiplicities up to 35 within a single sector",
        ("Not a defect in the HSF code — a property of the spectrum", 0, True),
        "options: restrict statistics to simple eigenvalues, or use a "
        "basis-independent quantity (eigenspace-averaged purity)",
        "worth checking which of the other plotted rules share this",
    ], notes="The mean stored entropy is sub-thermal throughout, 0.50–0.61 of "
             "ln2·⌊N/2⌋ — the direction the rest of the analysis predicts.")

    # 17 -----------------------------------------------------------------
    _bullets(prs, "Result 5 — why the same rule gives area AND volume law", [
        ("The three Fig. 3a states are in three very differently sized sectors",
         0, True),
        "single excitation |0…010…0⟩:  w = 2   →  |S_w| = C(N+1,2), POLYNOMIAL",
        "Néel |0101…⟩:                 w = N   →  |S_w| = N+1, LINEAR",
        "pair |001100…⟩:               w ≈ N/2 →  |S_w| exponential",
        ("So the premise fails: C(N+1,w) is exponentially large only near w = (N+1)/2",
         0, True),
        "the Néel state is a ONE-HOLE state in wall language — a wall on every "
        "bond but one, hard-core ⇒ JAMMED",
        ("The exact ceiling is the bipartite Schmidt rank, not ln|S_w|", 0, True),
        "from w = w_A + b_cut + w_B; the LOWER limit matters at high filling",
        "verified = the compressed Schmidt dimension for all w, N = 6…14",
        "at N=26: single ln 92 = 4.5, Néel ln 14 = 2.6, pair 9.011 = ⌊N/2⌋ln2 exactly",
    ], notes="R8 §9. I initially dropped the lower limit and caught it against "
             "the actual Schmidt dimensions — without it the Néel state looks "
             "unconstrained (64 patterns instead of 7 at N=12).")

    # 18 -----------------------------------------------------------------
    _picture(prs, "Area law, area law, volume law — predicted by w alone",
             os.path.join(FIGDIR, "fig_c150_entanglement.png"),
             "Left: S(t) at N=14 with the kinematic ceilings. Middle: the "
             "plateau vs N — single excitation and Néel are FLAT to N=38, the "
             "pair state grows. Right: the entanglement dome vs wall filling.",
             notes="Plateaus: single 1.00→1.24 and Néel 0.53→0.58 over N=10…38, "
                   "both flat and in fact flatter than their own ln N ceilings; "
                   "pair 1.29→2.45 over N=8…14 with a ceiling equal to the "
                   "volume-law value.")

    # 19 -----------------------------------------------------------------
    _bullets(prs, "Is the sector enough? — the three answers", [
        ("For the growth TYPE and its N-scaling: YES", 0, True),
        "via the exact ceiling, not via |S_w| — knowing w already separates ln N from N",
        ("For PRODUCT initial states, the sector nearly fixes the value too: ±3 %",
         0, True),
        "six random basis states of the N=14, w=6 shell: plateaus 2.26 … 2.39",
        "Fig. 3a uses product states ⇒ its three curves are set by three integers w",
        ("Over ALL states of a shell: NO", 0, True),
        "the stationary states of the same shell — C(⌈N/2⌉, w/2) = 35 at N=14, w=6 —",
        "have S(t) CONSTANT to 1e-15, at a value ABOVE the product plateau (3.5 vs 2.4)",
        "a 35-dimensional family with zero entanglement growth inside a volume-law sector",
        ("The dome's shape is the ceiling's shape; the dynamics supplies only "
         "an O(1) fill factor ≈ ½", 0, True),
        "the same sub-thermal fraction as the HSF entropies — as the degenerate "
        "spectrum predicts",
    ], notes="One sentence: the sector fixes the ceiling and hence the law; "
             "the dynamics fixes the prefactor; and only non-product states "
             "inside a sector can defy the sector's expectation, by not moving "
             "at all.")

    # 20 -----------------------------------------------------------------
    _bullets(prs, "Open items", [
        ("Identify the commutant", 0, True),
        "the 3^m law forces an exponentially large commutant (Σ m_λ²)",
        "its generators would prove the law AND make the block-resolved spacing "
        "statistics exact rather than a lower bound",
        ("Derive the two closed forms", 0, True),
        "dim Fix(U_w) = C(⌈N/2⌉, w/2), and a w-resolved version of the 3^m law",
        ("The partner-shell intertwiner", 0, True),
        "for odd N the shells w and N+1−w have identical spectra (1.2e-14),",
        "but the bond particle-hole map does NOT commute with U — unidentified",
        ("C150 on the ring", 0, True),
        "kinematics solved; the quantum layer is obc0 only. The 2-to-1 wall map "
        "gives the bond frame a Z₂ gauge redundancy",
        ("Frontier", 0, True),
        "N = 33 needs a 64-bit label array (68.7 GB), or a two-pass compact relabelling",
    ], notes="Ranked roughly by payoff. The commutant is the one that unlocks "
             "the others.")

    # 21 -----------------------------------------------------------------
    _bullets(prs, "Reproduce", [
        ("Closed forms, wall bijection, charge conservation, engine cross-check", 0, True),
        "python -m qca_fragmentation.c150 --verify",
        ("Frontiers", 0, True),
        "python -m qca_fragmentation.run_rule --rule 150 --N 21-22 --bc obc0",
        "python -m qca_fragmentation.c150 --frontier 21-32 --bc obc0",
        ("Quantum layer, level statistics, entanglement, HSF check", 0, True),
        "python -m qca_fragmentation.scaling.c150_report --quantum",
        "python -m qca_fragmentation.quantum.rule150_levels --run",
        "python -m qca_fragmentation.quantum.rule150_entanglement --run",
        "python -m qca_fragmentation.quantum.hsf_compare --all",
        ("Report and data", 0, True),
        "reports/pdf/R8_c150.pdf  (18 pp)   ·   analytics/c150_*.json",
        "full test suite: 403 tests",
    ], notes="This slide exists so the deck is self-contained if someone wants "
             "to rerun a number during the discussion.", font_size=16)

    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    return out


def main(argv=None):
    ap = argparse.ArgumentParser(description="C150 discussion deck")
    ap.add_argument("--out", default=OUT)
    args = ap.parse_args(argv)
    p = build(args.out)
    print("wrote", p)


if __name__ == "__main__":
    main()
